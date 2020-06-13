from .exportable import Exportable
from .company import Company
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor






class Companies(Exportable):
    companies = []

    def append(self, company):
        self.companies.append(company)

    def build(self):
        return Company()

    def write_pptx(self, name='default', argv=[]):

        # Create presentation (A4)
        pptx = Presentation()
        pptx.slide_height = PPTX_SIZE_A4_HEIGHT
        pptx.slide_width = PPTX_SIZE_A4_WIDTH
        title_slide_layout = pptx.slide_layouts[PPTX_LAYOUT_BLANK]

        # Add slide
        slide = pptx.slides.add_slide(title_slide_layout)

        # Slide: Title
        tit = slide.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(9.5), Inches(0.3))
        titp = tit.text_frame.paragraphs[0].text = argv.name
        titp = tit.text_frame.paragraphs[0].font.size = Pt(24)
        titp = tit.text_frame.paragraphs[0].font.bold = True

        # Slide: Subtitle
        subt = slide.shapes.add_textbox(Inches(0.2), Inches(0.6), Inches(9.5), Inches(0.3))
        subt.text_frame.paragraphs[0].text = "Subtitle lorem ipsum"
        subt.text_frame.paragraphs[0].font.size = Pt(20)
        subt.text_frame.paragraphs[0].font.color.rgb = PPTX_COLOR_DELOITTE_GREEN
        subt.text_frame.paragraphs[0].font.bold = True


        # Slide: Business info
        bizinft = slide.shapes.add_textbox(Inches(4.5), Inches(1), Inches(3), Inches(0.3))
        bizinft.text_frame.paragraphs[0].text = '企業概要 / Business overview'
        bizinft.text_frame.paragraphs[0].font.bold = True

        # Slide: Business info
        bizinft = slide.shapes.add_textbox(Inches(4.5), Inches(1.4), Inches(3), Inches(0.3))
        bizinft.text_frame.paragraphs[0].text = argv.name

        # Slide: Business info
        bizinft = slide.shapes.add_textbox(Inches(4.5), Inches(3), Inches(3), Inches(0.3))
        bizinft.text_frame.paragraphs[0].text = '製品概要 / Product overview'
        bizinft.text_frame.paragraphs[0].font.bold = True
        
        # Slide: Business info
        bizinft = slide.shapes.add_textbox(Inches(4.5), Inches(3.4), Inches(3), Inches(0.3))
        bizinft.text_frame.paragraphs[0].text = argv.name

        # Slide: Business info
        bizinft = slide.shapes.add_textbox(Inches(4.5), Inches(5), Inches(3), Inches(0.3))
        bizinft.text_frame.paragraphs[0].text = '協業仮設 / Collaboration hypothesis'
        bizinft.text_frame.paragraphs[0].font.bold = True
        
        # Slide: Business info
        bizinft = slide.shapes.add_textbox(Inches(4.5), Inches(5.4), Inches(3), Inches(0.3))
        bizinft.text_frame.paragraphs[0].text = argv.name
        bizinft.text_frame.paragraphs[0].level = 1

        # Slide: Company info
        company_info = slide.shapes.add_table(10, 2, Inches(0.2), Inches(1), Inches(4), Inches(3)).table
        company_info.columns[0].width = Inches(1.0)
        company_info.columns[1].width = Inches(3.0)
        company_info.cell(0, 0).text = '基本情報 / Basic info'
        company_info.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(12)
        company_info.cell(0, 1).text = ''
        company_info.cell(1, 0).text = '会社名'
        company_info.cell(2, 0).text = '現在地'
        company_info.cell(3, 0).text = '設立'
        company_info.cell(4, 0).text = '従業員数'
        company_info.cell(5, 0).text = '調達ラウンド'
        company_info.cell(6, 0).text = '調達資金金額'
        company_info.cell(7, 0).text = 'ウェブサイト'
        company_info.cell(1, 1).text = str(argv.name)
        company_info.cell(2, 1).text = str(argv.location)
        company_info.cell(3, 1).text = str(argv.founded)
        company_info.cell(4, 1).text = str(argv.employees)
        company_info.cell(7, 1).text = str(argv.url)

        # Business info
        business_info = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(9.5), Inches(1))

        pptx.save(name + '.pptx')