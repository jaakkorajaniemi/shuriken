import re
import copy
import six

from openpyxl import load_workbook

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
PPTX_LAYOUT_TITLE = 0
PPTX_LAYOUT_TITLEANDCONTENT = 1
PPTX_LAYOUT_SECTHEADER = 2
PPTX_LAYOUT_TWOCONTENT = 3
PPTX_LAYOUT_COMPARISON = 4
PPTX_LAYOUT_TITLEONLY = 5
PPTX_LAYOUT_BLANK = 6
PPTX_LAYOUT_CONTENTCAPTION = 7
PPTX_LAYOUT_PICTURECAPTION = 8
PPTX_SIZE_A4_WIDTH = 10689336
PPTX_SIZE_A4_HEIGHT = 7562088
PPTX_COLOR_DGREEN = RGBColor(0x86, 0xBC, 0x25)
PPTX_COLOR_DBLACK = RGBColor(0x00, 0x00, 0x00)
PPTX_COLOR_DWHITE = RGBColor(0xFF, 0xFF, 0xFF)


class Exportable():
    datatype = None
    dataname = None
    objects = []

    # Bindings
    bindings = {}
    
    def bind(self, var, key):
        self.bindings[var] = key

    # Constructor
    def __init__(self, datatype='xlsx', dataname=None):
        self.datatype = datatype
        self.dataname = dataname

    # Importers
    def read(self, datatype=None, dataname=None):
        if datatype == None: datatype = self.datatype
        if dataname == None: dataname = self.dataname

        if datatype == 'xlsx': self.read_xlsx(dataname)
        if datatype == 'pptx': self.read_pptx(dataname)

    def read_xlsx(self, dataname, keep_bindings=False):
        wb = load_workbook(dataname)
        ws = wb.active
        print("Loaded workbook " + dataname)

        # Header mappings
        var_list = list(self.bindings.keys()) 
        val_list = list(self.bindings.values())
        mappings = {}

        for col in ws.iter_cols(min_row=1, max_col=15, max_row=1):
            for cell in col:
                if cell.value in self.bindings.values():
                    mappings[var_list[val_list.index(cell.value)]] = cell.column
                    #print(cell.value + " => " + var_list[val_list.index(cell.value)] + "; column key " + str(cell.column))

        print("Finished header mappings")

        # Object creation loop
        self.objects = []
        for row in ws.iter_rows(min_row=2):
            
            obj = self.build()
            for col in mappings:
                #print("Setting attribute " + str(col) + " to " + str(row[mappings[col]-1].value))
                setattr(obj, col, row[mappings[col]-1].value)
                
            self.objects.append(obj)

        print("Successfully imported objects from " + dataname)

        if not keep_bindings:
            self.bindings = {}

    def build(self):
        print("Object builder is not defined; could not build object")
        
    def read_pptx(self, dataname):
        pass

    # Export
    def write(self, datatype=None, dataname=None, template=None):
        if datatype == None: datatype = self.datatype
        if dataname == None: dataname = self.dataname

        if datatype == 'pptx':
            if template is not None:
                self.load_pptx(dataname=dataname, template=template)
    
    def load_pptx(self, dataname=None, keep_bindings=False, template=None):
        print("Loading PowerPoint template " + template + " for output.")

        pptx = Presentation(template)
        text_runs = []

        # Header mappings
        var_list = list(self.bindings.keys()) 
        val_list = list(self.bindings.values())
        mappings = []

        for slide_id, slide in enumerate(pptx.slides):
            for shape_id, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                for paragraph_id, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_id, run in enumerate(paragraph.runs):
                        print(paragraph.text)

                        match = re.search("{{(.*)}}", paragraph.text)
                        if match: 
                            if match.group(1).strip() in self.bindings.values():
                                mappings.append([var_list[val_list.index(match.group(1).strip())], slide_id, shape_id, paragraph_id, run_id, match.group(0)])

        print("Finished header mappings")
        print(mappings)

        # Write a slide according to the template for each object in repository
        for obj_id, obj in enumerate(self.objects):
            slide = self.duplicate_slide(pptx, 0)

            for map in mappings:
                string = pptx.slides[map[1]].shapes[map[2]].text_frame.paragraphs[map[3]].text
                string = string.replace(map[5], str(getattr(obj, map[0], 'N/A')))
                font = copy.deepcopy(pptx.slides[map[1]].shapes[map[2]].text_frame.paragraphs[map[3]].runs[map[4]].font)
                pptx.slides[obj_id+1].shapes[map[2]].text_frame.paragraphs[map[3]].text = string
                pptx.slides[obj_id+1].shapes[map[2]].text_frame.paragraphs[map[3]].font.size = font.size
                pptx.slides[obj_id+1].shapes[map[2]].text_frame.paragraphs[map[3]].font.bold = font.bold
                if font.color.type:
                    pptx.slides[obj_id+1].shapes[map[2]].text_frame.paragraphs[map[3]].font.color.theme_color = font.color.theme_color

        pptx.save(dataname)
        print("Saved output to output")

        if not keep_bindings:
            self.bindings = {}

    def duplicate_slide(self, pres, index):
        template = pres.slides[index]
        try:
            blank_slide_layout = pres.slide_layouts[PPTX_LAYOUT_BLANK]
        except:
            blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

        copied_slide = pres.slides.add_slide(blank_slide_layout)

        for shp in template.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for _, value in six.iteritems(template.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                copied_slide.part.rels.add_relationship(
                    value.reltype,
                    value._target,
                    value.rId
                )

        return copied_slide